from flask import Flask, render_template, request, redirect, url_for, session, send_from_directory, send_file, jsonify
import requests
from youtube_transcript_api import YouTubeTranscriptApi
import google.generativeai as genai
import os
import shutil
from pathlib import Path
import time
from duckduckgo_search import DDGS
import re
from pytube import YouTube
from werkzeug.utils import secure_filename
import PyPDF2
from pptx import Presentation
import yt_dlp

genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# temp_dir = Path("temp_images")
temp_dir = Path(os.getenv("TEMP", "/tmp")) / "temp_images"


try:
    shutil.rmtree(temp_dir)
except FileNotFoundError:
    pass
except Exception as e:
    print(f"Warning: Could not delete temp_images - {e}")

temp_dir.mkdir(exist_ok=True)

PROMPT = """
You are a YouTube video summarizer designed for Mumbai University engineering students. 
Your task is to summarize the video transcript following the proper answer-writing format for 8-10 mark questions.

### **Instructions for Summarization:** 1. **Definition:** Start with a definition of the main topic and any closely related concepts.  
2. **Classification:** If the topic is broad, provide a **classification in a tree format** (use text-based representation like code blocks if needed).  
3. **Explanation:** Explain the topic in a structured, **stepwise or pointwise manner** to ensure clarity.  
4. **Diagrams:** If a diagram is necessary, Mention **"Draw a ____ Type of Diagram"** 5. **Merits & Demerits:** List advantages and disadvantages **if applicable**.  
6. **Applications:** Mention real-world applications **if applicable**.    
7. **Conclusion:** End with a brief 2-3 line conclusion summarizing the key points.  
"""

CONCISE_PROMPT = """
You are a YouTube video summarizer. Create a concise summary of the video in 5-10 key points.

Guidelines:
1. Each point should be clear and concise (1-2 lines max)
2. Use bullet points (•)
3. Focus on the most important concepts/ideas
4. Use keywords and technical terms where relevant
5. Keep the total summary within 200 words
6. Make points easy to remember and understand

Please provide a concise summary of this transcript:
"""

PDF_PPT_PROMPT = """
You are an educational content summarizer designed for engineering students. Analyze the provided content and create a comprehensive yet concise summary following this structure:

1. **Chapter Overview:**
   - Main topic and its significance
   - Key concepts covered
   - Prerequisites needed

2. **Topics Breakdown:**
   - List main topics and subtopics
   - Show relationships between concepts
   - Highlight important terms/definitions

3. **Simplified Explanations:**
   - Break down complex concepts
   - Use simple language
   - Provide examples where possible

4. **Key Points Summary:**
   - Bullet points of crucial information
   - Important formulas/equations (if any)
   - Common applications

5. **Study Focus:**
   - What to concentrate on
   - Potential exam topics
   - Common misconceptions to avoid

6. **Quick Revision Notes:**
   - 5-6 most important takeaways
   - Critical formulas/concepts to remember
   - Practice suggestion areas

Please analyze and summarize the following content:
"""

app = Flask(__name__)
app.secret_key = os.urandom(24)


@app.template_filter('format_content')
def format_content(text):
    if '|' in text:
        lines = text.split('\n')
        table_lines = []
        in_table = False
        formatted_text = []

        for line in lines:
            if line.strip().startswith('|'):
                if not in_table:
                    in_table = True
                    table_lines = ['<div class="table-responsive"><table class="comparison-table">']
                cells = line.strip().split('|')[1:-1]
                if '-|-' in line:
                    continue
                row = '<tr>' + ''.join(f'<td>{cell.strip()}</td>' for cell in cells) + '</tr>'
                table_lines.append(row)
            else:
                if in_table:
                    in_table = False
                    table_lines.append('</table></div>')
                    formatted_text.append('\n'.join(table_lines))
                formatted_text.append(line)

        if in_table:
            table_lines.append('</table></div>')
            formatted_text.append('\n'.join(table_lines))

        text = '\n'.join(formatted_text)

    text = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', text)

    return text


@app.route('/', methods=['GET', 'POST'])
def home():
    if request.method == 'POST':
        youtube_link = request.form.get('youtube_link')
        prompt_option = request.form.get('prompt_option')
        custom_prompt = request.form.get('custom_prompt', '')
        download_media_option = request.form.get('download_media') == 'on'  # get download option

        if youtube_link:
            session['youtube_link'] = youtube_link
            session['prompt_option'] = prompt_option
            session['custom_prompt'] = custom_prompt
            session['download_media'] = download_media_option  # add to session
            return redirect(url_for('process_video'))
    return render_template('index.html')


@app.route('/process_video')
def process_video():
    youtube_link = session.get('youtube_link')
    prompt_option = session.get('prompt_option')
    custom_prompt = session.get('custom_prompt', '')
    download_media_option = session.get('download_media')  # get from session

    if not youtube_link:
        return "Error: No YouTube link provided."

    video_id = youtube_link.split("v=")[-1].split("&")[0]
    thumbnail_url = f"http://img.youtube.com/vi/{video_id}/0.jpg"

    transcript = extract_transcript(video_id)

    if "Error" in transcript:
        return render_template('result.html',
                               thumbnail_url=thumbnail_url,
                               summary=transcript,
                               images=[],
                               transcript="")

    if prompt_option == "concise":
        summary = generate_summary(transcript, CONCISE_PROMPT)
        images = []
    elif prompt_option == "custom":
        summary = generate_summary(transcript, custom_prompt)
        images = []
    else:
        summary = generate_summary(transcript, PROMPT)
        diagram_query = determine_diagram_query(summary)
        images = download_relevant_images(diagram_query)
        if isinstance(images, str):
            images = []

    if "Error" in summary:
        return render_template('result.html',
                               thumbnail_url=thumbnail_url,
                               summary=summary,
                               images=[],
                               transcript="")

    session['summary'] = summary
    session['images'] = images
    session['current_image_index'] = 0

    return render_template('result.html',
                           thumbnail_url=thumbnail_url,
                           summary=summary,
                           images=images,
                           transcript=transcript,
                           youtube_link=youtube_link,
                           download_media=download_media_option)  # pass download option


@app.route('/static/temp_images/<path:filename>')
def serve_image(filename):
    return send_from_directory('temp_images', filename)


@app.route('/download_media/<video_id>/<media_type>')
def download_media(video_id, media_type):
    try:
        url = f"https://www.youtube.com/watch?v={video_id}"
        downloads_dir = Path('downloads')
        downloads_dir.mkdir(exist_ok=True)

        if media_type in ['mp4', 'mp3']:
            try:
                output_template = str(downloads_dir / f'%(title)s_{int(time.time())}')

                ydl_opts = {
                    'format': 'best' if media_type == 'mp4' else 'bestaudio/best',
                    'outtmpl': output_template,
                    'quiet': True,
                    'no_warnings': True,
                }

                if media_type == 'mp3':
                    ydl_opts.update({
                        'postprocessors': [{
                            'key': 'FFmpegExtractAudio',
                            'preferredcodec': 'mp3',
                            'preferredquality': '192',
                        }],
                    })

                with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                    info = ydl.extract_info(url, download=True)
                    downloaded_file = ydl.prepare_filename(info)

                    if media_type == 'mp3':
                        downloaded_file = Path(downloaded_file).with_suffix('.mp3')

                    if not Path(downloaded_file).exists():
                        return f"Failed to download {media_type}", 500

                    try:
                        return send_file(
                            str(downloaded_file),
                            as_attachment=True,
                            download_name=f"{info['title']}.{media_type}",
                            mimetype=f'video/mp4' if media_type == 'mp4' else 'audio/mp3'
                        )
                    finally:
                        try:
                            if Path(downloaded_file).exists():
                                os.remove(downloaded_file)
                        except Exception as e:
                            print(f"Cleanup error: {e}")

            except Exception as e:
                print(f"{media_type.upper()} download error: {str(e)}")
                return f"{media_type.upper()} download failed: {str(e)}", 500

        elif media_type == 'thumbnail':
            try:
                thumbnail_urls = [
                    f"https://img.youtube.com/vi/{video_id}/maxresdefault.jpg",
                    f"https://img.youtube.com/vi/{video_id}/hqdefault.jpg",
                    f"https://img.youtube.com/vi/{video_id}/mqdefault.jpg",
                    f"https://img.youtube.com/vi/{video_id}/default.jpg"
                ]

                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                }

                for thumb_url in thumbnail_urls:
                    response = requests.get(thumb_url, headers=headers)
                    if response.status_code == 200:
                        filename = f"thumbnail_{video_id}_{int(time.time())}.jpg"
                        thumbnail_path = downloads_dir / filename

                        thumbnail_path.write_bytes(response.content)

                        if not thumbnail_path.exists():
                            continue

                        try:
                            return send_file(
                                str(thumbnail_path),
                                as_attachment=True,
                                download_name=f"thumbnail_{video_id}.jpg",
                                mimetype='image/jpeg'
                            )
                        finally:
                            try:
                                if thumbnail_path.exists():
                                    os.remove(thumbnail_path)
                            except Exception as e:
                                print(f"Cleanup error: {e}")

                return "Could not download thumbnail", 400

            except Exception as e:
                print(f"Thumbnail download error: {str(e)}")
                return f"Thumbnail download failed: {str(e)}", 500

        else:
            return "Invalid media type specified", 400

    except Exception as e:
        print(f"Download error: {str(e)}")
        return f"Download failed: {str(e)}", 500


def extract_transcript(video_id):
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['en', 'hi'])
        full_transcript = ' '.join([entry['text'] for entry in transcript])
        return full_transcript
    except Exception as e:
        return f"Error extracting transcript: {str(e)}"


def generate_summary(transcript, custom_prompt=""):
    try:
        model = genai.GenerativeModel("gemini-1.5-flash")
        prompt_to_use = (custom_prompt + transcript) if custom_prompt else (PROMPT + transcript)
        response = model.generate_content(prompt_to_use)

        if hasattr(response, "text"):
            return response.text
        else:
            return "Error: Unexpected response format from API."
    except Exception as e:
        return f"Error generating summary: {str(e)}"


def determine_diagram_query(summary):
    if "8086" in summary:
        return "8086 microprocessor architecture block diagram"
    else:
        keywords = summary.split()[:5]
        return " ".join(keywords) + " diagram"


def download_relevant_images(search_query):
    try:
        with DDGS() as ddgs:
            image_results = ddgs.images(search_query, max_results=5)
        if not image_results:
            return []

        downloaded_images = []
        for index, result in enumerate(image_results[:5]):
            image_url = result["image"]
            try:
                response = requests.get(image_url, stream=True, timeout=10)
                if response.status_code == 200:
                    filename = f"image_{index}.jpg"
                    image_path = temp_dir / filename
                    with open(image_path, "wb") as f:
                        for chunk in response.iter_content(1024):
                            f.write(chunk)
                    downloaded_images.append(filename)
            except Exception as e:
                print(f"Error downloading image {index}: {e}")
                continue

        return downloaded_images
    except Exception as e:
        return f"Error downloading images: {str(e)}"


@app.route('/process_youtube', methods=['POST'])
def process_youtube():
    youtube_link = request.form.get('youtube_link')
    prompt_option = request.form.get('prompt_option')
    custom_prompt = request.form.get('custom_prompt', '')
    download_media_option = request.form.get('download_media') == 'on'  # get download option

    if youtube_link:
        session['youtube_link'] = youtube_link
        session['prompt_option'] = prompt_option
        session['custom_prompt'] = custom_prompt
        session['download_media'] = download_media_option  # add to session
        return redirect(url_for('process_video'))
    return redirect(url_for('home'))


@app.route('/process_document', methods=['POST'])
def process_document():
    if 'document' not in request.files:
        return "No file uploaded"

    file = request.files['document']
    if file.filename == '':
        return "No file selected"

    try:
        filename = secure_filename(file.filename)
        file_path = os.path.join('uploads', filename)

        os.makedirs('uploads', exist_ok=True)

        file.save(file_path)

        if filename.endswith('.pdf'):
            text = extract_pdf_text(file_path)
        elif filename.endswith(('.ppt', '.pptx')):
            text = extract_ppt_text(file_path)
        else:
            return "Unsupported file format"

        if not text.strip():
            return "No text could be extracted from the document"

        summary = generate_summary(text, PDF_PPT_PROMPT)

        os.remove(file_path)

        return render_template('document_result.html',
                               summary=summary,
                               text=text)

    except Exception as e:
        print(f"Error processing document: {str(e)}")
        return f"Error processing document: {str(e)}"


def extract_pdf_text(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text


def extract_ppt_text(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


@app.route('/ask_question', methods=['POST'])
def ask_question():
    try:
        data = request.json
        question = data.get('question')
        context = data.get('context')
        summary = data.get('summary')

        if not question or not context:
            return jsonify({
                'error': 'Missing question or context'
            }), 400

        prompt = f"""
        Based on the following document content and its summary, please answer the question.
        If the answer cannot be found in the content, say so.

        Document Summary:
        {summary}

        Full Document Content:
        {context}

        Question: {question}

        Please provide a clear and concise answer, using information from the document.
        If the information is not in the document, respond with "I cannot find information about this in the document."
        """

        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)

        if not response.text:
            return jsonify({
                'error': 'No response generated'
            }), 500

        return jsonify({
            'answer': response.text
        })

    except Exception as e:
        print(f"Error in ask_question: {str(e)}")
        return jsonify({
            'error': str(e)
        }), 500


if __name__ == '__main__':
    app.run(debug=True)